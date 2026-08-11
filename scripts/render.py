#!/usr/bin/env python3
"""Rendering helpers: turn a plain content spec into real PDF / PPTX / DOCX files.

Each document module in docs_content/ exposes DOC = {...} and this module turns it
into a file on disk. Keeping rendering separate from content means the corpus can be
regenerated or reformatted without touching the finance text.

Content spec (a dict):
  filename   'FIN-ACC-014_Accruals_Policy.pdf'
  title      document title
  subtitle   optional strap line
  meta       list of (label, value) shown in a header table
  blocks     list of blocks, each one of:
               ('h1', text)            section heading
               ('h2', text)            sub heading
               ('p',  text)           paragraph
               ('bullets', [text...]) bulleted list
               ('table', [[hdr...], [row...], ...])
               ('callout', text)       boxed emphasis (policy rule / warning)
               ('slide', title, [bullet...])   PPTX only
"""
import os

from reportlab.lib import colors
from reportlab.lib.enums import TA_LEFT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import (KeepTogether, ListFlowable, ListItem, PageBreak,
                                Paragraph, SimpleDocTemplate, Spacer, Table,
                                TableStyle)

OUT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "docs_out")

NAVY = colors.HexColor("#0B2545")
SLATE = colors.HexColor("#4A5568")
ACCENT = colors.HexColor("#1B6CA8")
LIGHT = colors.HexColor("#EDF2F7")
RULE = colors.HexColor("#CBD5E0")


# --------------------------------------------------------------------- styles
def _styles():
    ss = getSampleStyleSheet()
    s = {}
    s["title"] = ParagraphStyle("t", parent=ss["Title"], fontName="Helvetica-Bold",
                                fontSize=19, leading=23, textColor=NAVY,
                                alignment=TA_LEFT, spaceAfter=2)
    s["subtitle"] = ParagraphStyle("st", parent=ss["Normal"], fontName="Helvetica",
                                   fontSize=10.5, leading=14, textColor=SLATE,
                                   spaceAfter=10)
    s["h1"] = ParagraphStyle("h1", parent=ss["Heading1"], fontName="Helvetica-Bold",
                             fontSize=12.5, leading=16, textColor=NAVY,
                             spaceBefore=13, spaceAfter=5)
    s["h2"] = ParagraphStyle("h2", parent=ss["Heading2"], fontName="Helvetica-Bold",
                             fontSize=10.8, leading=14, textColor=ACCENT,
                             spaceBefore=9, spaceAfter=3)
    s["p"] = ParagraphStyle("p", parent=ss["BodyText"], fontName="Helvetica",
                            fontSize=9.6, leading=14, textColor=colors.HexColor("#1A202C"),
                            spaceAfter=6)
    s["li"] = ParagraphStyle("li", parent=s["p"], spaceAfter=3)
    s["callout"] = ParagraphStyle("c", parent=s["p"], fontName="Helvetica-Bold",
                                  fontSize=9.6, leading=14, textColor=NAVY)
    s["cell"] = ParagraphStyle("cell", parent=s["p"], fontSize=8.6, leading=11.5,
                               spaceAfter=0)
    s["cellh"] = ParagraphStyle("cellh", parent=s["cell"], fontName="Helvetica-Bold",
                                textColor=colors.white)
    return s


def _table(rows, st, col_widths=None):
    """Build a styled table; first row is treated as the header."""
    data = [[Paragraph(str(c), st["cellh"] if r == 0 else st["cell"])
             for c in row] for r, row in enumerate(rows)]
    t = Table(data, colWidths=col_widths, repeatRows=1, hAlign="LEFT")
    t.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), NAVY),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, LIGHT]),
        ("GRID", (0, 0), (-1, -1), 0.4, RULE),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 5),
        ("RIGHTPADDING", (0, 0), (-1, -1), 5),
        ("TOPPADDING", (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
    ]))
    return t


def render_pdf(doc: dict) -> str:
    """Render a content spec to a PDF and return the path."""
    os.makedirs(OUT_DIR, exist_ok=True)
    path = os.path.join(OUT_DIR, doc["filename"])
    st = _styles()
    story = []

    story.append(Paragraph(doc["title"], st["title"]))
    if doc.get("subtitle"):
        story.append(Paragraph(doc["subtitle"], st["subtitle"]))

    if doc.get("meta"):
        story.append(_table([["Field", "Detail"]] + [list(m) for m in doc["meta"]],
                            st, col_widths=[5.0 * cm, 11.5 * cm]))
        story.append(Spacer(1, 8))

    for block in doc["blocks"]:
        kind = block[0]
        if kind == "h1":
            story.append(Paragraph(block[1], st["h1"]))
        elif kind == "h2":
            story.append(Paragraph(block[1], st["h2"]))
        elif kind == "p":
            story.append(Paragraph(block[1], st["p"]))
        elif kind == "bullets":
            story.append(ListFlowable(
                [ListItem(Paragraph(b, st["li"]), leftIndent=12) for b in block[1]],
                bulletType="bullet", bulletFontSize=7, leftIndent=14, bulletOffsetY=1))
            story.append(Spacer(1, 4))
        elif kind == "table":
            story.append(Spacer(1, 2))
            story.append(_table(block[1], st))
            story.append(Spacer(1, 8))
        elif kind == "callout":
            box = Table([[Paragraph(block[1], st["callout"])]],
                        colWidths=[16.5 * cm], hAlign="LEFT")
            box.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, -1), LIGHT),
                ("BOX", (0, 0), (-1, -1), 0.8, ACCENT),
                ("LEFTPADDING", (0, 0), (-1, -1), 9),
                ("RIGHTPADDING", (0, 0), (-1, -1), 9),
                ("TOPPADDING", (0, 0), (-1, -1), 7),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 7),
            ]))
            story.append(KeepTogether(box))
            story.append(Spacer(1, 8))
        elif kind == "pagebreak":
            story.append(PageBreak())

    def _footer(canv, docp):
        canv.saveState()
        canv.setFont("Helvetica", 7.5)
        canv.setFillColor(SLATE)
        canv.drawString(2 * cm, 1.25 * cm,
                        f"{doc.get('footer', 'Meridian Bank - Group Finance')}  |  "
                        f"{doc.get('classification', 'Internal')}")
        canv.drawRightString(A4[0] - 2 * cm, 1.25 * cm, f"Page {docp.page}")
        canv.setStrokeColor(RULE)
        canv.line(2 * cm, 1.65 * cm, A4[0] - 2 * cm, 1.65 * cm)
        canv.restoreState()

    SimpleDocTemplate(
        path, pagesize=A4,
        leftMargin=2 * cm, rightMargin=2 * cm,
        topMargin=1.8 * cm, bottomMargin=2.1 * cm,
        title=doc["title"], author="Meridian Bank Group Finance",
        subject=doc.get("subtitle", ""),
    ).build(story, onFirstPage=_footer, onLaterPages=_footer)
    return path


def render_pptx(doc: dict) -> str:
    """Render a content spec to a PPTX deck and return the path."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    os.makedirs(OUT_DIR, exist_ok=True)
    path = os.path.join(OUT_DIR, doc["filename"])
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    def _txt(tf, text, size, bold=False, rgb=(26, 32, 44)):
        tf.text = text
        p = tf.paragraphs[0]
        p.font.size, p.font.bold = Pt(size), bold
        p.font.color.rgb = RGBColor(*rgb)
        p.font.name = "Calibri"

    # title slide
    s = prs.slides.add_slide(prs.slide_layouts[6])
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.4))
    _txt(box.text_frame, doc["title"], 40, True, (11, 37, 69))
    if doc.get("subtitle"):
        b2 = s.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.9))
        _txt(b2.text_frame, doc["subtitle"], 18, False, (74, 85, 104))
    if doc.get("meta"):
        b3 = s.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.4))
        tf = b3.text_frame
        for i, (k, v) in enumerate(doc["meta"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"{k}: {v}"
            p.font.size, p.font.name = Pt(12), "Calibri"
            p.font.color.rgb = RGBColor(74, 85, 104)

    # content slides
    for block in doc["blocks"]:
        if block[0] == "slide":
            _, stitle, bullets = block
            s = prs.slides.add_slide(prs.slide_layouts[6])
            bar = s.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12.0), Inches(0.9))
            _txt(bar.text_frame, stitle, 26, True, (11, 37, 69))
            body = s.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(11.6), Inches(5.3))
            tf = body.text_frame
            tf.word_wrap = True
            for i, b in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                sub = b.startswith("  ")
                p.text = ("- " if sub else "• ") + b.strip()
                p.font.size = Pt(13 if sub else 15)
                p.font.name = "Calibri"
                p.font.color.rgb = RGBColor(74, 85, 104) if sub else RGBColor(26, 32, 44)
                p.level = 1 if sub else 0
                p.space_after = Pt(7)
        elif block[0] == "slide_img":
            _, stitle, bullets, img_path = block
            s = prs.slides.add_slide(prs.slide_layouts[6])
            bar = s.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12.0), Inches(0.9))
            _txt(bar.text_frame, stitle, 26, True, (11, 37, 69))

            # Check if image exists; if not, fall back to bullets-only
            if img_path and os.path.exists(img_path):
                if bullets:
                    # Bullets on left ~55%, image on right ~45%
                    body = s.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(6.8), Inches(5.3))
                    tf = body.text_frame
                    tf.word_wrap = True
                    for i, b in enumerate(bullets):
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        sub = b.startswith("  ")
                        p.text = ("- " if sub else "• ") + b.strip()
                        p.font.size = Pt(12 if sub else 14)
                        p.font.name = "Calibri"
                        p.font.color.rgb = RGBColor(74, 85, 104) if sub else RGBColor(26, 32, 44)
                        p.level = 1 if sub else 0
                        p.space_after = Pt(5)
                    # Image on the right
                    s.shapes.add_picture(img_path, Inches(7.8), Inches(1.6),
                                        width=Inches(5.0))
                else:
                    # Full-width image below title
                    s.shapes.add_picture(img_path, Inches(0.9), Inches(1.6),
                                        width=Inches(11.6))
            else:
                # Fallback: bullets only
                body = s.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(11.6), Inches(5.3))
                tf = body.text_frame
                tf.word_wrap = True
                for i, b in enumerate(bullets):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    sub = b.startswith("  ")
                    p.text = ("- " if sub else "• ") + b.strip()
                    p.font.size = Pt(13 if sub else 15)
                    p.font.name = "Calibri"
                    p.font.color.rgb = RGBColor(74, 85, 104) if sub else RGBColor(26, 32, 44)
                    p.level = 1 if sub else 0
                    p.space_after = Pt(7)

    prs.save(path)
    return path


def render_docx(doc: dict) -> str:
    """Render a content spec to a DOCX and return the path."""
    from docx import Document
    from docx.shared import Pt, RGBColor

    os.makedirs(OUT_DIR, exist_ok=True)
    path = os.path.join(OUT_DIR, doc["filename"])
    d = Document()
    st = d.styles["Normal"]
    st.font.name, st.font.size = "Calibri", Pt(10.5)

    h = d.add_heading(doc["title"], level=0)
    for r in h.runs:
        r.font.color.rgb = RGBColor(0x0B, 0x25, 0x45)
    if doc.get("subtitle"):
        p = d.add_paragraph(doc["subtitle"])
        p.runs[0].italic = True

    if doc.get("meta"):
        t = d.add_table(rows=0, cols=2)
        t.style = "Light Grid Accent 1"
        for k, v in doc["meta"]:
            c = t.add_row().cells
            c[0].text, c[1].text = str(k), str(v)
        d.add_paragraph()

    for block in doc["blocks"]:
        kind = block[0]
        if kind == "h1":
            d.add_heading(block[1], level=1)
        elif kind == "h2":
            d.add_heading(block[1], level=2)
        elif kind == "p":
            d.add_paragraph(_plain(block[1]))
        elif kind == "bullets":
            for b in block[1]:
                d.add_paragraph(_plain(b), style="List Bullet")
        elif kind == "table":
            rows = block[1]
            t = d.add_table(rows=1, cols=len(rows[0]))
            t.style = "Light Grid Accent 1"
            for i, c in enumerate(rows[0]):
                t.rows[0].cells[i].text = _plain(str(c))
            for row in rows[1:]:
                cells = t.add_row().cells
                for i, c in enumerate(row):
                    cells[i].text = _plain(str(c))
            d.add_paragraph()
        elif kind == "callout":
            p = d.add_paragraph()
            p.add_run(_plain(block[1])).bold = True

    d.save(path)
    return path


def _plain(text: str) -> str:
    """Strip the light HTML used for PDF emphasis, for DOCX/PPTX output."""
    for a, b in (("<b>", ""), ("</b>", ""), ("<i>", ""), ("</i>", ""),
                 ("<br/>", " "), ("&amp;", "&")):
        text = text.replace(a, b)
    return text


RENDERERS = {".pdf": render_pdf, ".pptx": render_pptx, ".docx": render_docx}


def render(doc: dict) -> str:
    ext = os.path.splitext(doc["filename"])[1].lower()
    return RENDERERS[ext](doc)
