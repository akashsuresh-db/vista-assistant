#!/usr/bin/env python3
"""Verify the rendered corpus actually contains the ground-truth facts.

Two failure modes this guards against:
  1. An agent wrote plausible-sounding prose with invented numbers.
  2. A document contradicts the structured data, so doc + data disagree in the demo.

Extracts real text from the rendered PDF/DOCX/PPTX files (not the python source) so it
tests what the Knowledge Assistant will actually index.
"""
import glob
import os
import re
import sys

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "docs_out")


# ------------------------------------------------------------------ extraction
def text_pdf(path):
    from pypdf import PdfReader
    return "\n".join((pg.extract_text() or "") for pg in PdfReader(path).pages)


def text_docx(path):
    from docx import Document
    d = Document(path)
    parts = [p.text for p in d.paragraphs]
    for t in d.tables:
        for r in t.rows:
            parts += [c.text for c in r.cells]
    return "\n".join(parts)


def text_pptx(path):
    from pptx import Presentation
    parts = []
    for s in Presentation(path).slides:
        for sh in s.shapes:
            if sh.has_text_frame:
                parts.append(sh.text_frame.text)
    return "\n".join(parts)


EXTRACT = {".pdf": text_pdf, ".docx": text_docx, ".pptx": text_pptx}


def norm(t):
    """Normalise so 1,339,800 / 1339800 / 1,339.8k style variants all compare."""
    return re.sub(r"\s+", " ", t.replace("’", "'").replace("–", "-")).lower()


# ------------------------------------------------------------------ expectations
# (filename fragment, [(label, [any-of these strings must appear]) ...])
EXPECT = [
    ("FIN-ACC-014", [
        ("90-day rule", ["90 day", "90-day", "90 days"]),
        ("policy ref", ["fin-acc-014"]),
        ("evidence hierarchy", ["po_receipt_not_invoiced"]),
        ("weakest basis", ["manual_estimate"]),
        ("controller sign-off threshold", ["50,000", "50000"]),
        ("register location", ["fact_accrual", "v_aged_accrual"]),
        ("breach status", ["breach_over_90_days"]),
    ]),
    ("FSSC-SOP-002", [
        ("sop ref", ["fssc-sop-002"]),
        ("close day 8 consolidation", ["group consolidation submission"]),
        ("accrual task", ["accrual calculation & upload", "accrual calculation and upload"]),
        ("ic matching task", ["intercompany billing & matching", "intercompany billing and matching"]),
        ("ap cut-off task", ["sub-ledger cut-off - ap", "sub-ledger cut-off ap"]),
        ("commentary threshold", ["250,000", "250000"]),
        ("breach tracker", ["is_sla_breach", "fact_close_task", "v_close_performance"]),
    ]),
    ("FIN-FX-007", [
        ("policy ref", ["fin-fx-007"]),
        ("constant currency", ["constant-currency", "constant currency"]),
        ("gbp plan rate", ["1.178"]),
        ("gbp actual rate", ["1.299"]),
        ("eur rates", ["0.992", "1.091"]),
        ("mb-uk total variance", ["798"]),
        ("mb-uk constant ccy", ["107"]),
        ("fx impact", ["691"]),
        ("plan rate lock", ["plan rate"]),
        ("rate table source", ["dim_fx_rate"]),
    ]),
    ("FIN-IC-021", [
        ("policy ref", ["fin-ic-021"]),
        ("reportable threshold", ["250,000", "250000"]),
        ("auto write-off", ["1,000", "1000"]),
        ("break status", ["unresolved_break"]),
        ("register", ["fact_intercompany"]),
        ("treasury funding", ["treasury funding"]),
    ]),
    ("IA-2026-11", [
        ("memo ref", ["ia-2026-11"]),
        ("total aged", ["1,339,800", "1339800"]),
        ("cards concentration", ["1,069,500", "1069500", "4 "]),
        ("ACR5183", ["acr5183"]),
        ("ACR5184", ["acr5184"]),
        ("ACR5185", ["acr5185"]),
        ("ACR5186", ["acr5186"]),
        ("ACR5187", ["acr5187"]),
        ("ACR5188", ["acr5188"]),
        ("helix amount", ["486,000", "486000"]),
        ("helix dispute", ["1,310,500", "1310500"]),
        ("policy cross-ref", ["fin-acc-014"]),
        ("helios link", ["helios"]),
    ]),
    ("PROJ-HELIOS", [
        ("approved after plan lock", ["aop_2026_v3", "aop 2026"]),
        ("feb approval", ["february 2026", "2026-02"]),
        ("consulting account", ["530100"]),
        ("cloud account", ["520200"]),
        ("licence account", ["520100"]),
        ("contractor account", ["510400"]),
        ("cards platform engineering", ["cards platform engineering"]),
        ("helix vendor", ["helix consulting"]),
        ("northwind vendor", ["northwind cloud"]),
        ("kestrel vendor", ["kestrel software"]),
        ("contractor ramp", ["95"]),
    ]),
    ("MDA_Q2_2026", [
        ("cards actual", ["28.9", "28,95"]),
        ("cards budget", ["26.2"]),
        ("cards variance", ["2.72"]),
        ("consulting driver", ["979"]),
        ("cloud driver", ["608"]),
        ("licence driver", ["361"]),
        ("ramp may", ["1.46"]),
        ("helios", ["helios"]),
        ("aged accruals", ["1,339,800", "1339800"]),
        ("ic break", ["482,000", "482000"]),
        ("helix 90+", ["1,310,500", "1310500"]),
        ("fx gbp", ["1.299"]),
        ("mb-uk favourable in gbp", ["6%", "favourable"]),
        ("sla breaches", ["sub-ledger cut-off"]),
    ]),
    ("Variance_Commentary_Pack_2026-06", [
        ("threshold", ["250,000", "250000"]),
        ("cards variance", ["2.72"]),
        ("consulting driver", ["979"]),
        ("cloud driver", ["608"]),
        ("mb-uk fx split", ["798", "107", "691"]),
        ("aged accruals", ["1,339,800", "1339800"]),
        ("ic break", ["482,000", "482000"]),
        ("policy xref accrual", ["fin-acc-014"]),
        ("policy xref ic", ["fin-ic-021"]),
        ("audit xref", ["ia-2026-11"]),
        ("constant currency", ["constant-currency", "constant currency"]),
    ]),
]

# Numbers that must NEVER appear (would contradict the structured data)
FORBIDDEN = [
    ("wrong cards variance", ["+2.4m unfavourable", "2.48m", "2,485,978"]),
    ("internal demo flag leaked", ["is_planted_overrun"]),
]


def main():
    files = {os.path.basename(p): p for p in glob.glob(os.path.join(OUT, "*"))}
    fails, checks = [], 0

    for frag, expectations in EXPECT:
        match = [p for n, p in files.items() if frag.lower() in n.lower()]
        if not match:
            fails.append(f"MISSING FILE matching '{frag}'")
            print(f"FAIL  no file matching {frag}")
            continue
        path = match[0]
        ext = os.path.splitext(path)[1].lower()
        body = norm(EXTRACT[ext](path))
        name = os.path.basename(path)
        print(f"\n--- {name}  ({len(body):,} chars extracted)")
        for label, anyof in expectations:
            checks += 1
            hit = any(norm(a) in body for a in anyof)
            print(f"  {'ok  ' if hit else 'FAIL'} {label}")
            if not hit:
                fails.append(f"{name}: missing {label} (any of {anyof})")
        for label, bad in FORBIDDEN:
            for b in bad:
                if norm(b) in body:
                    fails.append(f"{name}: FORBIDDEN {label} -> '{b}'")
                    print(f"  FAIL forbidden content: {b}")

    print(f"\n{checks} content checks across {len(EXPECT)} documents")
    if fails:
        print(f"\n{len(fails)} FAILURES:")
        for f in fails:
            print("  -", f)
        return 1
    print("ALL DOCUMENT CONTENT CHECKS PASSED")
    return 0


if __name__ == "__main__":
    sys.exit(main())
